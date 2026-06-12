# =========================================================
# SAI'S AI UNIVERSE - COMPLETE ULTRA AI PLATFORM
# =========================================================

# =========================================================
# INSTALL THESE FIRST
# =========================================================

# pip install streamlit groq python-dotenv PyPDF2 gtts
# pip install streamlit-mic-recorder python-pptx python-docx
# pip install reportlab pillow requests speechrecognition
# pip install beautifulsoup4 lxml pandas youtube-transcript-api
# pip install pytz 

# =========================================================
# IMPORTS
# =========================================================

import streamlit as st
import streamlit.components.v1 as components
from groq import Groq
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from gtts import gTTS
from streamlit_mic_recorder import mic_recorder

from pptx import Presentation
from pptx.util import Inches

from docx import Document
from reportlab.pdfgen import canvas

from PIL import Image
from bs4 import BeautifulSoup

import speech_recognition as sr
import pandas as pd
from youtube_transcript_api import YouTubeTranscriptApi

import tempfile
import requests
import sqlite3
import hashlib
import os
import base64
import datetime
import urllib.parse
import pytz
import json
import time
import random
import io

# =========================================================
# PAGE CONFIG - UPDATED
# =========================================================

st.set_page_config(
    page_title="Sai's AI Universe",  
    layout="wide",
    initial_sidebar_state="expanded"
)

# =========================================================
# CUSTOM CSS (TWINKLING STARS & ORBITRON FONT)
# =========================================================

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Orbitron:wght@400;600;800&family=Inter:wght@300;400;600;700;800&display=swap');

/* --- Twinkling Stars Background --- */
.stApp {
    background-color: #010204;
    background-image: 
        radial-gradient(white, rgba(255,255,255,.1) 1.5px, transparent 2px), 
        radial-gradient(white, rgba(255,255,255,.07) 1px, transparent 1.5px), 
        radial-gradient(rgba(108, 99, 255, 0.05) 15%, transparent 60%); 
    background-size: 550px 550px, 350px 350px, 100% 100%;
    background-position: 0 0, 40px 60px, 0 0;
    color: #FFFFFF;
    font-family: 'Inter', sans-serif;
    perspective: 2000px; overflow-x: hidden;
    animation: realGalaxyMove 60s linear infinite;
}

.stApp::before {
    content: ""; position: fixed; top: 0; left: 0; width: 100%; height: 100%;
    background: radial-gradient(circle at 15% 25%, rgba(108, 99, 255, 0.15), transparent 35%),
                radial-gradient(circle at 85% 75%, rgba(138, 43, 226, 0.12), transparent 35%);
    z-index: -2;
    animation: nebulaTwist 120s ease-in-out infinite;
}

.stApp::after {
    content: ""; position: fixed; top: 0; left: 0; width: 100%; height: 100%;
    background-image: radial-gradient(white, rgba(255,255,255,.2) 2px, transparent 3px);
    background-size: 650px 650px;
    z-index: -1;
    animation: twinkling 4s ease-in-out infinite; opacity: 0.8;
}

@keyframes realGalaxyMove { from { background-position: 0 0, 40px 60px, 0 0; } to { background-position: 550px 550px, 390px 410px, 0 0; } }
@keyframes nebulaTwist { 0% { transform: rotate(0deg); } 50% { transform: rotate(2deg) scale(1.05); } 100% { transform: rotate(0deg); } }
@keyframes twinkling { 0% { opacity: 0.3; } 50% { opacity: 1.0; } 100% { opacity: 0.3; } }

/* Fonts for Headings */
h1, h2, h3 { font-family: 'Orbitron', sans-serif !important; color: #FFFFFF !important; }
h4, h5, h6, p, label { color: #E2E8F0 !important; }

[data-testid="stSidebar"] {
    background-color: rgba(14, 18, 29, 0.8) !important;
    backdrop-filter: blur(15px);
    border-right: 1px solid rgba(255,255,255,0.05) !important;
}

.tool-card {
    background: rgba(22, 27, 42, 0.6); 
    backdrop-filter: blur(10px); 
    border: 1px solid rgba(255,255,255,0.1);
    border-radius: 16px; 
    padding: 25px; 
    text-align: left; 
    transition: all 0.4s cubic-bezier(0.175, 0.885, 0.32, 1.275);
    cursor: pointer; 
    height: 100%;
    transform-style: preserve-3d;
    perspective: 1000px;
}
.tool-card:hover {
    border-color: #6C63FF; 
    transform: translateY(-10px) rotateX(6deg) rotateY(-4deg) scale(1.03); 
    box-shadow: -10px 15px 30px rgba(0, 0, 0, 0.5), 0 0 20px rgba(108, 99, 255, 0.4);
}

.login-left {
    background: rgba(14, 16, 26, 0.9); backdrop-filter: blur(10px); padding: 60px;
    border-radius: 24px; border: 1px solid #2A2E3F;
}

.stButton button {
    background: linear-gradient(135deg, #6C63FF 0%, #8B5CF6 100%) !important;
    border-radius: 10px !important;
    border: none !important;
    font-weight: 600 !important;
    color: white !important;
    padding: 10px 20px !important;
    width: 100%;
}
.stTextInput input, .stTextArea textarea, .stSelectbox > div > div {
    background-color: rgba(22, 27, 42, 0.7) !important;
    border: 1px solid rgba(255,255,255,0.1) !important;
    border-radius: 10px !important;
    color: white !important;
}
</style>
""", unsafe_allow_html=True)

# =========================================================
# LOAD ENV & SECRETS (WORKS FOR BOTH LOCAL & LIVE APP)
# =========================================================

load_dotenv()

try:
    api_key = st.secrets.get("GROQ_API_KEY", os.getenv("GROQ_API_KEY"))
    cf_api_token = st.secrets.get("CLOUDFLARE_API_TOKEN", os.getenv("CLOUDFLARE_API_TOKEN", "")).strip().strip('"').strip("'")
    cf_account_id = st.secrets.get("CLOUDFLARE_ACCOUNT_ID", os.getenv("CLOUDFLARE_ACCOUNT_ID", "")).strip().strip('"').strip("'")
except:
    api_key = os.getenv("GROQ_API_KEY")
    cf_api_token = os.getenv("CLOUDFLARE_API_TOKEN", "").strip().strip('"').strip("'")
    cf_account_id = os.getenv("CLOUDFLARE_ACCOUNT_ID", "").strip().strip('"').strip("'")

if not api_key:
    st.error("Please Add GROQ_API_KEY In .env File or Streamlit Cloud Secrets")
    st.stop()

# =========================================================
# CONNECT GROQ
# =========================================================

client = Groq(api_key=api_key)

# =========================================================
# DATABASE
# =========================================================

conn = sqlite3.connect(
    "users.db",
    check_same_thread=False
)

cursor = conn.cursor()

cursor.execute("""
CREATE TABLE IF NOT EXISTS users (
    username TEXT,
    password TEXT
)
""")

try:
    cursor.execute("ALTER TABLE users ADD COLUMN email TEXT")
except:
    pass

cursor.execute("""
CREATE TABLE IF NOT EXISTS vault (
    username TEXT, 
    filename TEXT, 
    filedata BLOB
)
""")

# AUTO-LOGIN SESSIONS TABLE
cursor.execute("""
CREATE TABLE IF NOT EXISTS sessions (
    token TEXT,
    username TEXT
)
""")

conn.commit()

# =========================================================
# PASSWORD HASH
# =========================================================

def hash_password(password):
    return hashlib.sha256(
        password.encode()
    ).hexdigest()

# =========================================================
# WORLDWIDE LIVE SEARCH 
# =========================================================

def get_live_search_data(query):
    context = ""
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
    }

    from datetime import datetime, timedelta
    ist = pytz.timezone('Asia/Kolkata')
    
    now_utc = datetime.now(pytz.utc)
    now_ist = now_utc.astimezone(ist)
    
    today = now_ist.date()
    yesterday = today - timedelta(days=1)
    
    search_query = query.lower()
    if "yesterday" in search_query or "yesterdays" in search_query:
        search_query = search_query.replace("yesterdays", yesterday.strftime("%B %d, %Y")).replace("yesterday", yesterday.strftime("%B %d, %Y"))
    if "today" in search_query or "todays" in search_query:
        search_query = search_query.replace("todays", today.strftime("%B %d, %Y")).replace("today", today.strftime("%B %d, %Y"))
        
    try:
        url = "https://lite.duckduckgo.com/lite/"
        data = {"q": search_query}
        res = requests.post(url, headers=headers, data=data, timeout=6)
        if res.status_code == 200:
            soup = BeautifulSoup(res.text, 'html.parser')
            snippets = soup.find_all('td', class_='result-snippet')
            for snippet in snippets[:4]:
                context += f"- {snippet.text.strip()}\n"
            if context.strip():
                return context
    except:
        pass

    try:
        google_url = f"https://www.google.com/search?q={urllib.parse.quote(search_query)}"
        res = requests.get(google_url, headers=headers, timeout=6)
        soup = BeautifulSoup(res.text, 'html.parser')
        for div in soup.find_all('div', class_=['BNeawe', 's3v9rd', 'AP7Wnd'])[:5]:
            text = div.get_text()
            if len(text) > 20:
                context += f"- {text.strip()}\n"
        if context.strip():
            return context
    except:
        pass

    try:
        wiki_url = f"https://en.wikipedia.org/w/api.php?action=query&list=search&srsearch={urllib.parse.quote(search_query)}&utf8=&format=json"
        wiki_res = requests.get(wiki_url, timeout=5).json()
        if 'query' in wiki_res and 'search' in wiki_res['query']:
            for item in wiki_res['query']['search'][:3]:
                clean_text = BeautifulSoup(item['snippet'], 'html.parser').text
                context += f"- {item['title']}: {clean_text}\n"
    except:
        pass

    return context

# =========================================================
# STATE MANAGEMENT FOR MENU 
# =========================================================

if "menu" not in st.session_state:
    st.session_state.menu = "🏠 Dashboard"

if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
    
if "username" not in st.session_state:
    st.session_state.username = None

# =========================================================
# PURE & CLEAN LOGIN SYSTEM (NO GOOGLE / NO COOKIES)
# =========================================================

if not st.session_state.logged_in:

    # --- NORMAL LOGIN UI ---
    st.markdown("<br><br>", unsafe_allow_html=True)
    l_col, r_col, _ = st.columns([1.2, 1, 0.2])
    
    with l_col:
        st.markdown("""
        <div class="login-left">
            <h1 style="font-size: 3.2rem; line-height: 1.2;">Welcome to <br><span style="color: #6C63FF;">Sai's AI Universe</span></h1>
            <p style="color: #A0A4B8; font-size: 1.2rem; margin-top: 15px;">Intelligent, secure, and built to help you achieve more every single day.</p>
        </div>
        """, unsafe_allow_html=True)

    with r_col:
        st.markdown("<h2>Welcome back</h2><p style='color: #A0A4B8;'>Sign in to your account to continue</p>", unsafe_allow_html=True)

        auth_option = st.radio("Choose Option", ["Login", "Signup"], horizontal=True)

        username = st.text_input("Username")
        email = st.text_input("Email")
        password = st.text_input("Password", type="password")

        if auth_option == "Signup":
            if st.button("Create Account"):
                hashed = hash_password(password)
                cursor.execute("INSERT INTO users (username, email, password) VALUES (?, ?, ?)", (username, email, hashed))
                conn.commit()
                st.success("Account Created Successfully")

        else:
            if st.button("Login"):
                hashed = hash_password(password)
                cursor.execute("SELECT * FROM users WHERE username=? AND password=?", (username, hashed))
                result = cursor.fetchone()

                if result:
                    st.session_state.logged_in = True
                    st.session_state.username = username
                    
                    st.success("Login Successful!")
                    time.sleep(1)
                    st.rerun()
                else:
                    st.error("Invalid Username Or Password")

# =========================================================
# MAIN APP (DASHBOARD)
# =========================================================

else:

    # =====================================================
    # MODERN SIDEBAR
    # =====================================================
    
    st.sidebar.title("⚡ Sai's AI Universe")

    col1, col2 = st.sidebar.columns([1, 3])
    with col1:
        st.write("👤") 
    with col2:
        st.write(f"**{st.session_state.username}**")

    st.sidebar.markdown("---")
    
    if st.sidebar.button("🏠 Home Dashboard", use_container_width=True):
        st.session_state.menu = "🏠 Dashboard"
        st.rerun()

    if st.sidebar.button("⚙️ Settings", use_container_width=True):
        st.session_state.menu = "⚙️ Settings"
        st.rerun()
        
    st.sidebar.markdown("---")

    model = st.sidebar.selectbox("Choose Model", ["llama-3.3-70b-versatile", "mixtral-8x7b-32768", "gemma2-9b-it"])
    temperature = st.sidebar.slider("Creativity", 0.0, 1.0, 0.7)

    # =====================================================
    # MEMORY & AI FUNCTION (UPGRADED FOR MAX ACCURACY & NO CRASHES)
    # =====================================================

    if "messages" not in st.session_state:
        st.session_state.messages = []

    def generate_response(prompt):
        current_tool = st.session_state.menu
        messages_list = []
        
        # MASTER PROMPTS INJECTED HERE SO WE DON'T TOUCH YOUR LOGIC BELOW
        master_prompts = {
            "💻 Coding Assistant": "Act as an Elite Principal Software Engineer. Write clean, perfectly optimized, robust, and bug-free code. Include deep explanations.",
            "📄 PDF Analyzer": "You are a Master Document Analyst. Analyze the text carefully and answer strictly based on the text provided. Give an exhaustive and detailed response.",
            "🧠 Study Assistant": "You are a brilliant professor. Explain complex topics using simple analogies and structured formats. Ensure 100% factual accuracy.",
            "🌐 Translator": "You are an expert multi-lingual translator. Translate the text perfectly while maintaining the original tone and intent.",
            "📝 Resume Builder": "You are an expert HR Manager. Format details into a highly professional, ATS-friendly resume structure.",
            "📚 Quiz Generator": "You are a master educator. Generate a structured MCQ quiz with exact answers and deep explanations.",
            "📧 Email Writer": "You are an expert in corporate communication. Write a professional, polite, and highly effective email.",
            "📰 Article Writer": "You are a Pulitzer Prize-winning journalist. Write a captivating, well-researched, and highly detailed article.",
            "📱 Caption Generator": "You are an expert Social Media Manager. Create highly engaging, viral captions with relevant hashtags.",
            "🧮 Math Solver": "You are an expert Mathematician. Solve the problem step-by-step with 100% logical accuracy. Check your calculations.",
            "📊 PPT Generator": "You are a master presentation designer. Provide brilliant slide-by-slide content.",
            "📄 PDF Generator": "You are a professional technical writer. Write highly detailed, well-structured, and accurate notes.",
            "📝 DOCX Generator": "You are a professional author. Write a comprehensive, perfectly structured document.",
            "🌐 Web Summarizer": "You are an expert data analyst. Extract the core arguments, facts, and structure a perfect, detailed summary.",
            "📊 Data Analyzer": "You are a Lead Data Scientist. Analyze the dataset strictly based on the stats provided. Give extremely accurate, mathematical, and logical answers without hallucinating numbers.",
            "🎬 YT Summarizer": "You are an expert content analyst. Break down this video transcript into key ideas, quotes, and a structured conclusion.",
            "🛠️ Code Reviewer": "Act as an Elite Security & Performance Code Reviewer. Find logical bugs, security flaws, and performance issues. Explain deeply, then provide the perfectly fixed code.",
            "💡 Idea Validator": "Act as a Top Tier Silicon Valley VC. Ruthlessly evaluate the idea. Give a massive, highly detailed report containing Market Size, Competitors, Monetization, SWOT, and a 30-Day Execution Roadmap."
        }
        
        if current_tool not in ["🏠 Dashboard", "💬 AI Chat"]:
            system_msg = f"You are 'Sai's Universe AI', a Master-Level Expert executing the '{current_tool}' tool.\n"
            system_msg += master_prompts.get(current_tool, "You are a STRICT and HIGHLY SPECIALIZED tool. Your ONLY purpose is to perform tasks related to this tool.")
            system_msg += "\nCRITICAL RULES:\n1. BE EXHAUSTIVE: Provide detailed, well-structured, and comprehensive responses.\n2. BE 100% ACCURATE: Your logic and facts must be perfect.\n3. CURRENT YEAR: It is currently 2026."
            messages_list.append({"role": "system", "content": system_msg})
            
        with st.spinner("Fetching latest updates from internet..."):
            live_context = get_live_search_data(prompt)
            if live_context and live_context.strip():
                prompt = prompt + f"\n\n--- STRICT LATEST WEB DATA FOR ACCURACY ---\n{live_context}"
                
        messages_list.append({"role": "user", "content": prompt})
        
        try:
            response = client.chat.completions.create(
                model=model,
                messages=messages_list,
                temperature=0.2 if current_tool != "💬 AI Chat" else temperature,
                max_tokens=6000 # Added 6000 max_tokens here to prevent app crash on big inputs
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"⚠️ API Error (Please try splitting the text if it's too large): {e}"

    # =====================================================
    # ADVANCED DASHBOARD GRID UI 
    # =====================================================

    menu = st.session_state.menu

    if menu == "🏠 Dashboard":
        
        ist = pytz.timezone('Asia/Kolkata')
        current_hour = datetime.datetime.now(ist).hour
        greeting = "Good morning" if current_hour < 12 else "Good afternoon" if current_hour < 18 else "Good evening"

        st.markdown(f"<h1>{greeting}, {st.session_state.username} 👋</h1>", unsafe_allow_html=True)
        
        search_col, filter_col = st.columns([3, 1])
        with search_col:
            search_query = st.text_input("🔍 Search Tools...", placeholder="e.g., Video, Chat, Resume")
        with filter_col:
            category = st.selectbox("📁 Category", ["All", "AI Chat & Text", "Media & Generation", "Utility & Work"])
        
        st.write("")

        tools = [
            ("💬 AI Chat", "Chat with intelligent AI", "AI Chat & Text"),
            ("💻 Coding Assistant", "Write and debug code", "Utility & Work"),
            ("📄 PDF Analyzer", "Extract info from PDFs", "Utility & Work"),
            ("🎨 AI Image Generator", "Create stunning images", "Media & Generation"),
            ("🔊 Voice Generator", "Text to speech conversion", "Media & Generation"),
            ("🧠 Study Assistant", "Learn any topic easily", "AI Chat & Text"),
            ("🌐 Translator", "Translate between languages", "AI Chat & Text"),
            ("📝 Resume Builder", "Create professional resumes", "Utility & Work"),
            ("📚 Quiz Generator", "Generate MCQ quizzes", "Utility & Work"),
            ("📧 Email Writer", "Draft professional emails", "Utility & Work"),
            ("📰 Article Writer", "Write detailed articles", "Utility & Work"),
            ("📱 Caption Generator", "Social media captions", "Media & Generation"),
            ("🧮 Math Solver", "Step-by-step math solutions", "Utility & Work"),
            ("📊 PPT Generator", "Create AI presentations", "Media & Generation"),
            ("📄 PDF Generator", "Create document PDFs", "Utility & Work"),
            ("📝 DOCX Generator", "Create Word documents", "Utility & Work"),
            ("🖼 Image To PDF", "Convert photos to PDF", "Utility & Work"),
            ("🤖 Advanced AI Agents", "Multi-role AI assistants", "AI Chat & Text"),
            ("📁 My Vault", "Save and manage your files", "Utility & Work"),
            ("👁️ Vision AI", "Analyze images with AI", "Media & Generation"),
            ("🌐 Web Summarizer", "Summarize webpage links", "Utility & Work"),
            ("🎙️ Audio To Text", "Convert audio to text", "Media & Generation"),
            ("📊 Data Analyzer", "Analyze CSV datasets", "Utility & Work"),
            ("🎬 YT Summarizer", "Summarize YouTube videos", "Media & Generation"),
            ("🛠️ Code Reviewer", "Fix and review code", "Utility & Work"),
            ("💡 Idea Validator", "Validate business ideas", "Utility & Work")
        ]

        filtered_tools = [t for t in tools if (search_query.lower() in t[0].lower() or search_query.lower() in t[1].lower()) and (category == "All" or t[2] == category)]

        for i in range(0, len(filtered_tools), 3):
            cols = st.columns(3)
            for j in range(3):
                if i + j < len(filtered_tools):
                    name, desc, cat = filtered_tools[i+j]
                    with cols[j]:
                        st.markdown(f"""
                        <div class="tool-card" style="box-shadow: 0 4px 15px rgba(108, 99, 255, 0.2);">
                            <h3 style="margin-bottom:5px; display:flex; justify-content:space-between;">{name}</h3>
                            <p style="color:#A0A4B8; font-size:0.9rem;">{desc}</p>
                            <span style="font-size:0.7rem; background:#23293B; padding:3px 8px; border-radius:10px;">{cat}</span>
                        </div>
                        """, unsafe_allow_html=True)
                        if st.button(f"Open", key=name, use_container_width=True):
                            st.session_state.menu = name
                            st.rerun()

    elif menu == "⚙️ Settings":
        st.title("⚙️ Account Settings")
        st.write("Manage your AI Universe preferences here.")
        st.info("Dark/Light mode is controlled by your system/browser settings in Streamlit natively.")
        st.toggle("Enable Notifications", value=True)
        st.toggle("Save Chat History to Vault")
        
    else:
        col1, col2 = st.columns([1, 5])
        with col1:
            if st.button("⬅ Back to Home"):
                st.session_state.menu = "🏠 Dashboard"
                st.rerun()
        st.markdown("---")

        if menu == "💬 AI Chat":

            st.title("🤖 AI Chat Assistant")

            for msg in st.session_state.messages:
                with st.chat_message(msg["role"]):
                    if isinstance(msg["content"], list):
                        for item in msg["content"]:
                            if item["type"] == "text":
                                st.markdown(item["text"])
                            elif item["type"] == "image_url":
                                st.image(item["image_url"]["url"])
                    else:
                        st.markdown(msg["content"])

            audio = mic_recorder(
                start_prompt="🎤 Start Recording",
                stop_prompt="⏹ Stop Recording",
                key="mic"
            )

            action_col1, action_col2 = st.columns([15, 1])

            with action_col2:
                with st.popover("➕"):  
                    st.markdown("**🛠️ Tools & Attachments**")
                    use_google = st.toggle("🌐 Web Search")
                    st.divider()
                    
                    attach_opt = st.radio("Attachment:", ["❌ None", "📁 Upload File", "📸 Open Camera"])
                    
                    attached_text = ""
                    attached_image_b64 = None
                    
                    if attach_opt == "📁 Upload File":
                        doc_file = st.file_uploader("Upload File", type=["png", "jpg", "jpeg", "txt", "pdf", "csv"])
                        if doc_file:
                            if doc_file.name.endswith((".png", ".jpg", ".jpeg")):
                                st.image(doc_file, width=200)
                                doc_file.seek(0)
                                attached_image_b64 = base64.b64encode(doc_file.read()).decode('utf-8')
                                st.success("Attached!")
                            elif doc_file.name.endswith(".txt"):
                                attached_text = f"\n\n[Attached File Content]:\n" + doc_file.read().decode("utf-8")
                                st.success("Attached!")
                            elif doc_file.name.endswith(".pdf"):
                                try:
                                    pdf_reader = PdfReader(doc_file)
                                    pdf_text = "".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
                                    attached_text = f"\n\n[Attached PDF Content]:\n" + pdf_text[:4000]
                                    st.success("Attached!")
                                except Exception as e:
                                    st.error(f"PDF Error: {e}")
                            elif doc_file.name.endswith(".csv"):
                                try:
                                    df_chat = pd.read_csv(doc_file)
                                    attached_text = f"\n\n[Attached CSV Data Sample]:\n" + df_chat.head(5).to_string()
                                    st.success("Attached!")
                                except Exception as e:
                                    st.error(f"CSV Error: {e}")
                                    
                    elif attach_opt == "📸 Open Camera":
                        st.info("💡 Camera is ON. Select 'None' to save battery.")
                        st.warning("📱 Note: Mobile camera will ONLY open if the app is running on HTTPS (deployed version) or localhost. It won't work on local Wi-Fi IPs.")
                        cam_file = st.camera_input("Take a photo")
                        if cam_file:
                            st.image(cam_file, width=200)
                            cam_file.seek(0)
                            attached_image_b64 = base64.b64encode(cam_file.read()).decode('utf-8')
                            st.success("Photo attached!")

            prompt = st.chat_input("Ask Anything...")

            if prompt:
                
                is_edit_request = attached_image_b64 and any(word in prompt.lower() for word in ["edit", "change", "make it", "filter", "black and white", "blur", "bright", "dark"])
                
                if is_edit_request:
                    user_message = {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{attached_image_b64}"}}
                        ]
                    }
                    st.session_state.messages.append(user_message)
                    
                    with st.chat_message("user"):
                        st.markdown(prompt)
                        st.caption("📸 Image Attached for Editing")
                        
                    with st.chat_message("assistant"):
                        st.markdown("🎨 Processing your image edit request...")
                        try:
                            from PIL import Image, ImageFilter, ImageEnhance
                            
                            img_data = base64.b64decode(attached_image_b64)
                            img = Image.open(io.BytesIO(img_data))
                            
                            if "black and white" in prompt.lower() or "grayscale" in prompt.lower():
                                img = img.convert("L")
                            elif "blur" in prompt.lower():
                                img = img.filter(ImageFilter.BLUR)
                            elif "bright" in prompt.lower():
                                enhancer = ImageEnhance.Brightness(img)
                                img = enhancer.enhance(1.5)
                            elif "dark" in prompt.lower():
                                enhancer = ImageEnhance.Brightness(img)
                                img = enhancer.enhance(0.5)
                            else:
                                st.warning("⚠️ For advanced edits (like changing backgrounds), paid APIs are needed. Applied Auto-Enhance instead!")
                                enhancer = ImageEnhance.Color(img)
                                img = enhancer.enhance(1.5)
                                
                            buf = io.BytesIO()
                            img.save(buf, format="PNG")
                            edited_b64 = base64.b64encode(buf.getvalue()).decode('utf-8')
                            
                            st.image(img, caption="Edited Result", use_container_width=True)
                            
                            st.session_state.messages.append(
                                {
                                    "role": "assistant",
                                    "content": [
                                        {"type": "text", "text": "Here is your edited image!"},
                                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{edited_b64}"}}
                                    ]
                                }
                            )
                        except Exception as e:
                            st.error(f"Failed to edit image: {e}")
                            
                else:
                    google_context = ""
                    if 'use_google' in locals() and use_google:
                        with st.spinner("Searching the Web Worldwide..."):
                            google_context = get_live_search_data(prompt)

                    final_content = prompt + attached_text

                    if attached_image_b64:
                        user_message = {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": final_content},
                                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{attached_image_b64}"}}
                            ]
                        }
                        chat_model = "llama-3.2-11b-vision-preview" 
                    else:
                        user_message = {
                            "role": "user",
                            "content": final_content
                        }
                        chat_model = model

                    st.session_state.messages.append(user_message)

                    with st.chat_message("user"):

                        st.markdown(prompt)
                        if attached_image_b64:
                            st.caption("📸 Image Attached")
                        if attached_text:
                            st.caption("📁 Document Attached")
                        if 'use_google' in locals() and use_google and google_context:
                            st.caption("🌐 Used Live Search Data")

                    with st.chat_message("assistant"):

                        placeholder = st.empty()
                        full_response = ""
                        
                        try:
                            ist = pytz.timezone('Asia/Kolkata')
                            current_date_str = datetime.datetime.now(ist).strftime("%B %d, %Y")
                            
                            user_name = st.session_state.username
                            system_prompt = f"You are a helpful AI assistant named Sai's Universe AI. You are chatting with {user_name}. Always refer to them by their name when appropriate and be friendly. Today's date is {current_date_str}. Note that the current year is 2026, and political positions have updated since the 2024 elections (e.g., Nara Chandrababu Naidu is the current active Chief Minister of Andhra Pradesh)."
                            
                            api_messages = [{"role": "system", "content": system_prompt}]
                            
                            for idx, m in enumerate(st.session_state.messages):
                                is_last_message = (idx == len(st.session_state.messages) - 1)
                                inject_text = ""
                                
                                if is_last_message and 'use_google' in locals() and use_google:
                                    if google_context.strip():
                                        inject_text = f"\n\n--- LATEST WEB DATA ---\n{google_context}\n\nCRITICAL RULE: The current year is 2026. Your internal database cutoff is outdated (2023). Do NOT use your old 2023 memories to guess current political leaders. For Andhra Pradesh, Nara Chandrababu Naidu is the Chief Minister since June 2024. Use the live web data and this 2026 timeframe to answer the user's query perfectly and naturally without any disclaimers."
                                    else:
                                        inject_text = f"\n\nCRITICAL RULE: Provide the best factual answer from your internal knowledge, but keep in mind that the year is 2026. DO NOT mention that you could not search the web. Just answer naturally."
                                        
                                if isinstance(m["content"], list):
                                    if chat_model != "llama-3.2-11b-vision-preview": 
                                        text_only = ""
                                        for item in m["content"]:
                                            if item["type"] == "text":
                                                text_only = item["text"]
                                        text_only += inject_text
                                        api_messages.append({"role": m["role"], "content": text_only})
                                    else:
                                        temp_m = m.copy()
                                        if inject_text:
                                            new_content = []
                                            for item in temp_m["content"]:
                                                if item["type"] == "text":
                                                    new_content.append({"type": "text", "text": item["text"] + inject_text})
                                                else:
                                                    new_content.append(item)
                                            temp_m["content"] = new_content
                                        api_messages.append(temp_m)
                                else:
                                    temp_content = m["content"]
                                    temp_content += inject_text
                                    api_messages.append({"role": m["role"], "content": temp_content})

                            response = client.chat.completions.create(
                                model=chat_model,
                                messages=api_messages,
                                temperature=0.2 if ('use_google' in locals() and use_google) else temperature, 
                                max_tokens=6000, # Added max_tokens here for massive inputs
                                stream=True
                            )

                            for chunk in response:
                                if chunk.choices[0].delta.content:
                                    content = chunk.choices[0].delta.content
                                    full_response += content
                                    placeholder.markdown(full_response + "▌")

                            placeholder.markdown(full_response)
                            
                        except Exception as e:
                            full_response = f"API Error: {e}"
                            placeholder.markdown(full_response)

                        st.session_state.messages.append(
                            {
                                "role": "assistant",
                                "content": full_response
                            }
                        )

        # =====================================================
        # CODING ASSISTANT
        # =====================================================

        elif menu == "💻 Coding Assistant":

            st.title("💻 Coding Assistant")

            code_prompt = st.text_area(
                "Enter Coding Question"
            )

            if st.button("Generate Code"):

                answer = generate_response(code_prompt)

                st.code(answer)

        # =====================================================
        # PDF ANALYZER
        # =====================================================

        elif menu == "📄 PDF Analyzer":

            st.title("📄 PDF Analyzer")

            uploaded_file = st.file_uploader(
                "Upload PDF",
                type="pdf"
            )

            if uploaded_file:

                pdf_reader = PdfReader(uploaded_file)

                text = "".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])

                st.success("PDF Loaded Successfully")

                question = st.text_input(
                    "Ask Question From PDF"
                )

                if st.button("Analyze PDF"):

                    prompt = f"""
                    PDF Content:
                    {text[:5000]}

                    Question:
                    {question}
                    """

                    answer = generate_response(prompt)

                    st.markdown(answer)

        # =====================================================
        # AI IMAGE GENERATOR (PERMANENT CLOUDFLARE FIX)
        # =====================================================

        elif menu == "🎨 AI Image Generator":

            st.title("🎨 AI Image Generator (Pro Mode)")
            st.write("Generate amazing high-quality images using your secure Cloudflare token!")

            if not cf_api_token or not cf_account_id:
                st.error("💡 Cloudflare API Token or Account ID Not Found! Please add 'CLOUDFLARE_API_TOKEN' and 'CLOUDFLARE_ACCOUNT_ID' to your Secrets or .env file.")
                st.stop()

            st.markdown("### ⚙️ Select Image Style")
            img_style = st.radio(
                "Choose your visual vibe:",
                [
                    "🌟 Ultra HD (Photorealistic)", 
                    "🎨 Anime / Manga Style",
                    "✨ Cinematic / 3D Render"
                ],
                horizontal=False
            )

            st.markdown("---")

            image_prompt = st.text_input(
                "Describe the image in detail:",
                placeholder="e.g., A cute robot reading a book in a magical forest"
            )

            if st.button("✨ Generate Image"):
                if image_prompt:
                    with st.spinner("🎨 Creating your masterpiece... Please wait..."):
                        try:
                            if "Ultra HD" in img_style:
                                final_prompt = image_prompt + ", highly detailed, photorealistic, 8k resolution, masterpiece"
                            elif "Anime" in img_style:
                                final_prompt = image_prompt + ", anime style, studio ghibli, vibrant colors, detailed illustration"
                            else:
                                final_prompt = image_prompt + ", cinematic lighting, octane render, unreal engine 5, highly detailed"
                            
                            API_URL = f"https://api.cloudflare.com/client/v4/accounts/{cf_account_id}/ai/run/@cf/stabilityai/stable-diffusion-xl-base-1.0"
                            headers = {"Authorization": f"Bearer {cf_api_token}"}
                            payload = {"prompt": final_prompt}
                            
                            response = requests.post(API_URL, headers=headers, json=payload, timeout=60)
                            
                            if response.status_code == 200:
                                img_bytes = response.content
                                img = Image.open(io.BytesIO(img_bytes))
                                
                                st.success("✅ Image Generated Successfully!")
                                st.image(img, caption=image_prompt, use_container_width=True)
                                
                                st.download_button(
                                    label="⬇️ Download High-Res Image",
                                    data=img_bytes,
                                    file_name="masterpiece.jpg",
                                    mime="image/jpeg",
                                    use_container_width=True
                                )
                            elif response.status_code == 503:
                                st.warning("⏳ The AI Model is currently waking up. Please wait 15 seconds and click Generate again!")
                            else:
                                try:
                                    error_msg = response.json()
                                except:
                                    error_msg = response.text
                                st.error(f"⚠️ Cloudflare Error: {error_msg}")
                                
                        except Exception as e:
                            st.error(f"Network Error: {e}. Check your internet connection.")
                else:
                    st.warning("Please describe what you want to generate!")

        # =====================================================
        # VOICE GENERATOR
        # =====================================================

        elif menu == "🔊 Voice Generator":

            st.title("🔊 Voice Generator")

            voice_text = st.text_area(
                "Enter Text"
            )

            if st.button("Generate Voice"):

                tts = gTTS(
                    text=voice_text,
                    lang='en'
                )

                audio_file = "temp_voice.mp3"

                tts.save(audio_file)

                st.audio(audio_file)

        # =====================================================
        # STUDY ASSISTANT
        # =====================================================

        elif menu == "🧠 Study Assistant":

            st.title("🧠 Study Assistant")

            topic = st.text_input(
                "Enter Topic"
            )

            if st.button("Explain Topic"):

                answer = generate_response(
                    f"Explain {topic} simply"
                )

                st.markdown(answer)

        # =====================================================
        # TRANSLATOR
        # =====================================================

        elif menu == "🌐 Translator":

            st.title("🌐 Translator")

            text = st.text_area("Enter Text")

            language = st.selectbox(
                "Translate To",
                [
                    "Telugu",
                    "Hindi",
                    "French",
                    "Spanish"
                ]
            )

            if st.button("Translate"):

                answer = generate_response(
                    f"Translate to {language}: {text}"
                )

                st.markdown(answer)

        # =====================================================
        # RESUME BUILDER
        # =====================================================

        elif menu == "📝 Resume Builder":

            st.title("📝 Resume Builder")

            details = st.text_area(
                "Enter Resume Details"
            )

            if st.button("Generate Resume"):

                answer = generate_response(
                    f"Create professional resume: {details}"
                )

                st.markdown(answer)

        # =====================================================
        # QUIZ GENERATOR
        # =====================================================

        elif menu == "📚 Quiz Generator":

            st.title("📚 Quiz Generator")

            topic = st.text_input("Quiz Topic")

            if st.button("Generate Quiz"):

                answer = generate_response(
                    f"Create MCQ quiz about {topic}"
                )

                st.markdown(answer)

        # =====================================================
        # EMAIL WRITER
        # =====================================================

        elif menu == "📧 Email Writer":

            st.title("📧 Email Writer")

            topic = st.text_input("Email Topic")

            if st.button("Generate Email"):

                answer = generate_response(
                    f"Write professional email about {topic}"
                )

                st.markdown(answer)

        # =====================================================
        # ARTICLE WRITER
        # =====================================================

        elif menu == "📰 Article Writer":

            st.title("📰 Article Writer")

            topic = st.text_input("Article Topic")

            if st.button("Generate Article"):

                answer = generate_response(
                    f"Write detailed article about {topic}"
                )

                st.markdown(answer)

        # =====================================================
        # CAPTION GENERATOR
        # =====================================================

        elif menu == "📱 Caption Generator":

            st.title("📱 Caption Generator")

            topic = st.text_input("Post Topic")

            if st.button("Generate Caption"):

                answer = generate_response(
                    f"Generate captions for {topic}"
                )

                st.markdown(answer)

        # =====================================================
        # MATH SOLVER
        # =====================================================

        elif menu == "🧮 Math Solver":

            st.title("🧮 Math Solver")

            problem = st.text_input(
                "Enter Math Problem"
            )

            if st.button("Solve"):

                answer = generate_response(
                    f"Solve step by step: {problem}"
                )

                st.markdown(answer)

                st.latex(
                    r"x = \frac{-b \pm \sqrt{b^2-4ac}}{2a}"
                )

        # =====================================================
        # PPT GENERATOR
        # =====================================================

        elif menu == "📊 PPT Generator":

            st.title("📊 AI PPT Generator")

            ppt_topic = st.text_input(
                "Enter PPT Topic"
            )

            slides_count = st.slider(
                "Slides",
                3,
                15,
                5
            )

            if st.button("Generate PPT"):

                prompt = f"""
                Create content for {slides_count} slides about:
                {ppt_topic}
                """

                answer = generate_response(prompt)

                prs = Presentation()

                slides_data = answer.split("Slide")

                for slide_data in slides_data:

                    if slide_data.strip():

                        slide_layout = prs.slide_layouts[6]

                        slide = prs.slides.add_slide(
                            slide_layout
                        )

                        title_box = slide.shapes.add_textbox(
                            Inches(0.5),
                            Inches(0.3),
                            Inches(8),
                            Inches(1)
                        )

                        title_box.text = ppt_topic

                        content_box = slide.shapes.add_textbox(
                            Inches(0.7),
                            Inches(1.5),
                            Inches(5),
                            Inches(3)
                        )

                        content_box.text = slide_data

                        headers = {
                            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
                            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
                            "Accept-Language": "en-US,en;q=0.5",
                            "Connection": "keep-alive",
                            "Upgrade-Insecure-Requests": "1"
                        }
                        
                        formatted_ppt_topic = urllib.parse.quote(ppt_topic)
                        seed = random.randint(1, 1000000)
                        image_url = f"https://image.pollinations.ai/prompt/{formatted_ppt_topic}?seed={seed}"
                            
                        try:
                            response = requests.get(image_url, headers=headers, timeout=20)
                            image_path = "ppt_image.jpg"

                            if response.status_code == 200 and len(response.content) > 1000:
                                img = Image.open(io.BytesIO(response.content))
                                img.save(image_path)
                                slide.shapes.add_picture(
                                    image_path,
                                    Inches(5.5),
                                    Inches(1.5),
                                    width=Inches(3)
                                )
                        except Exception:
                            pass

                ppt_file = f"{ppt_topic}.pptx"

                prs.save(ppt_file)

                with open(ppt_file, "rb") as file:

                    st.download_button(
                        label="⬇ Download PPT",
                        data=file,
                        file_name=ppt_file,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

        # =====================================================
        # PDF GENERATOR
        # =====================================================

        elif menu == "📄 PDF Generator":

            st.title("📄 PDF Generator")

            topic = st.text_input(
                "Enter PDF Topic"
            )

            if st.button("Generate PDF"):

                answer = generate_response(
                    f"Write notes about {topic}"
                )

                pdf_file = f"{topic}.pdf"

                c = canvas.Canvas(pdf_file)

                text_object = c.beginText(40, 800)

                for line in answer.split("\n"):

                    text_object.textLine(line)

                c.drawText(text_object)

                c.save()

                with open(pdf_file, "rb") as file:

                    st.download_button(
                        label="⬇ Download PDF",
                        data=file,
                        file_name=pdf_file,
                        mime="application/pdf"
                    )

        # =====================================================
        # DOCX GENERATOR
        # =====================================================

        elif menu == "📝 DOCX Generator":

            st.title("📝 DOCX Generator")

            topic = st.text_input(
                "Enter DOCX Topic"
            )

            if st.button("Generate DOCX"):

                answer = generate_response(
                    f"Write document about {topic}"
                )

                doc = Document()

                doc.add_heading(topic, level=1)

                doc.add_paragraph(answer)

                docx_file = f"{topic}.docx"

                doc.save(docx_file)

                with open(docx_file, "rb") as file:

                    st.download_button(
                        label="⬇ Download DOCX",
                        data=file,
                        file_name=docx_file,
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )

        # =====================================================
        # IMAGE TO PDF
        # =====================================================

        elif menu == "🖼 Image To PDF":

            st.title("🖼 Image To PDF")

            uploaded_image = st.file_uploader(
                "Upload Image",
                type=["png", "jpg", "jpeg"]
            )

            if uploaded_image:

                st.image(uploaded_image)

                if st.button("Convert To PDF"):

                    image = Image.open(uploaded_image)

                    pdf_file = "converted_image.pdf"

                    rgb_image = image.convert("RGB")

                    rgb_image.save(pdf_file)

                    with open(pdf_file, "rb") as file:

                        st.download_button(
                            label="⬇ Download PDF",
                            data=file,
                            file_name=pdf_file,
                            mime="application/pdf"
                        )

        # =====================================================
        # ADVANCED MULTI-ROLE AI AGENTS
        # =====================================================

        elif menu == "🤖 Advanced AI Agents":

            st.title("🤖 Advanced AI Agents")
            st.write("Select specialized AI agents for your specific tasks.")

            agent_type = st.selectbox(
                "Select Agent Persona",
                [
                    "🔬 Research Agent (Deep Web Research)",
                    "💻 Coding Agent (Software Architect)",
                    "📈 Finance Agent (Market Analysis)",
                    "🎓 Study Agent (Academic Tutor)",
                    "✈️ Travel Agent (Trip Planner)"
                ]
            )

            task = st.text_area(f"Give a task to the {agent_type.split(' ')[1]}:", height=150)

            if st.button(f"🚀 Run {agent_type.split(' ')[1]}"):
                if task:
                    with st.spinner(f"Agent {agent_type.split(' ')[1]} is processing..."):
                        
                        system_prompts = {
                            "🔬": "You are an expert Research Agent. Provide deep, factual, and well-structured research with citations where possible.",
                            "💻": "You are a Senior Software Architect. Provide robust, scalable, and bug-free code with clear explanations.",
                            "📈": "You are a Financial Analyst Agent. Provide logical market analysis, avoiding direct investment advice, focusing on trends and data.",
                            "🎓": "You are an Academic Study Agent. Explain concepts clearly, using analogies and step-by-step breakdowns for easy learning.",
                            "✈️": "You are a Luxury Travel Agent. Create detailed, practical, and exciting itineraries including logistics and local tips."
                        }
                        
                        icon = agent_type.split(" ")[0]
                        sys_prompt = system_prompts.get(icon, "You are a helpful AI.")

                        answer = generate_response(f"System: {sys_prompt}\n\nUser Task: {task}")
                        st.markdown(f"### 📋 Agent Report:\n{answer}")
                else:
                    st.warning("Please enter a task first.")


        # =====================================================
        # NEW FEATURES
        # =====================================================

        elif menu == "📁 My Vault":

            st.title("📁 My Vault")

            st.write("Securely store and retrieve your generated files.")

            v_col1, v_col2 = st.columns(2)

            with v_col1:

                st.markdown("### Upload to Vault")

                v_upload = st.file_uploader(
                    "Upload any file to Vault"
                )

                if st.button("Secure Save"):

                    if v_upload:

                        cursor.execute(
                            "INSERT INTO vault (username, filename, filedata) VALUES (?, ?, ?)", 
                            (st.session_state.username, v_upload.name, v_upload.read())
                        )

                        conn.commit()

                        st.success(f"{v_upload.name} securely saved!")

            with v_col2:

                st.markdown("### Your Saved Files")

                cursor.execute(
                    "SELECT filename, filedata FROM vault WHERE username=?", 
                    (st.session_state.username,)
                )

                saved_files = cursor.fetchall()

                if saved_files:

                    for fname, fdata in saved_files:

                        st.download_button(
                            label=f"⬇ Download {fname}", 
                            data=fdata, 
                            file_name=fname
                        )

                else:

                    st.info("Vault is empty.")

        elif menu == "👁️ Vision AI":

            st.title("👁️ Vision AI")

            img_file = st.file_uploader(
                "Upload Image for Analysis", 
                type=["png", "jpg", "jpeg"]
            )

            if img_file:

                st.image(img_file, width=300)

                img_prompt = st.text_input("Ask about this image:")

                if st.button("Analyze Image"):

                    img_file.seek(0)

                    base64_image = base64.b64encode(img_file.read()).decode('utf-8')

                    try:

                        vision_resp = client.chat.completions.create(
                            model="llama-3.2-11b-vision-preview", 
                            messages=[{
                                "role": "system",
                                "content": "You are a specialized Image Analysis AI. STRICT RULE: Answer questions ONLY about the provided image. If the user asks general chat questions ('hi', 'how are you') or asks you to do tasks unrelated to the image, reply EXACTLY with: '⚠️ Invalid request. I only analyze images.'"
                            },
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text": img_prompt or "Describe this image in detail."},
                                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                                ]
                            }],
                            temperature=0.3,
                            max_tokens=4000 # Added max_tokens for safety
                        )

                        st.success(vision_resp.choices[0].message.content)

                    except Exception as e:

                        st.error(f"Vision API Error. {e}")

        # =====================================================
        # WEB SUMMARIZER
        # =====================================================

        elif menu == "🌐 Web Summarizer":

            st.title("🌐 Web Summarizer")

            url = st.text_input("Enter Website URL")

            if st.button("Summarize"):

                try:

                    res = requests.get(url)
                    
                    soup = BeautifulSoup(res.text, 'html.parser')
                    
                    text = soup.get_text(separator=' ', strip=True)[:5000]

                    prompt = f"Summarize this webpage content in detail:\n\n{text}"
                    
                    answer = generate_response(prompt)
                    
                    st.markdown(answer)

                except Exception as e:
                    
                    st.error(f"Error fetching URL: {e}")

        # =====================================================
        # AUDIO TO TEXT
        # =====================================================

        elif menu == "🎙️ Audio To Text":

            st.title("🎙️ Audio To Text")

            st.write("Upload a .wav audio file to extract text.")

            audio_file = st.file_uploader("Upload Audio", type=["wav"])

            if audio_file is not None:

                st.audio(audio_file)

                if st.button("Extract Text"):

                    recognizer = sr.Recognizer()
                    
                    with sr.AudioFile(audio_file) as source:

                        audio_data = recognizer.record(source)

                        try:

                            text = recognizer.recognize_google(audio_data)
                            
                            st.success("Extraction Successful!")
                            
                            st.write(text)

                        except Exception as e:
                            
                            st.error(f"Error converting audio: {e}")

        # =====================================================
        # DATA ANALYZER
        # =====================================================

        elif menu == "📊 Data Analyzer":

            st.title("📊 Data Analyzer")

            csv_file = st.file_uploader("Upload CSV File", type=["csv"])

            if csv_file is not None:

                df = pd.read_csv(csv_file)
                
                st.dataframe(df.head())

                question = st.text_input("Ask a question about this data")

                if st.button("Analyze Data"):

                    data_sample = df.head(10).to_string()
                    
                    data_info = df.describe().to_string()
                    
                    prompt = f"Here is a sample of my dataset:\n{data_sample}\n\nHere is the summary info:\n{data_info}\n\nBased on this, answer the user's question: {question}"

                    answer = generate_response(prompt)
                    
                    st.markdown(answer)

        # =====================================================
        # YT SUMMARIZER
        # =====================================================

        elif menu == "🎬 YT Summarizer":

            st.title("🎬 YouTube Video Summarizer")

            yt_url = st.text_input("Enter YouTube Video Link")

            if st.button("Summarize Video"):

                try:

                    if "v=" in yt_url: 
                        video_id = yt_url.split("v=")[1].split("&")[0]
                    elif "youtu.be/" in yt_url:
                        video_id = yt_url.split("youtu.be/")[1].split("?")[0]
                    elif "shorts/" in yt_url:
                        video_id = yt_url.split("shorts/")[1].split("?")[0]
                    else:
                        st.error("Invalid YouTube URL")
                        video_id = None

                    if video_id:

                        try:
                            transcript_list = YouTubeTranscriptApi().fetch(video_id)
                        except AttributeError:
                            transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
                        
                        transcript_text = " ".join([t.text if hasattr(t, 'text') else t['text'] for t in transcript_list])[:5000]

                        prompt = f"Summarize this YouTube video transcript in detail:\n\n{transcript_text}"
                        
                        answer = generate_response(prompt)
                        
                        st.markdown(answer)

                except Exception as e:
                    
                    st.error(f"Could not fetch transcript. The video might not have subtitles. Error: {e}")

        # =====================================================
        # CODE REVIEWER
        # =====================================================

        elif menu == "🛠️ Code Reviewer":

            st.title("🛠️ Code Bug Fixer & Reviewer")

            bug_code = st.text_area("Paste your code here with errors")

            if st.button("Review & Fix Code"):

                prompt = f"Act as an expert software engineer. Review the following code, find any bugs, explain them, and provide the corrected code:\n\n{bug_code}"

                answer = generate_response(prompt)
                
                st.markdown(answer)

        # =====================================================
        # IDEA VALIDATOR
        # =====================================================

        elif menu == "💡 Idea Validator":

            st.title("💡 Startup Idea Validator")

            st.write("Turn your raw business ideas into a complete execution plan.")

            idea = st.text_area("Describe your startup or business idea:")

            if st.button("Validate Idea"):

                if idea:

                    prompt = f"""
                    Act as an expert Silicon Valley Startup Consultant and Investor.
                    Analyze this business idea: "{idea}"
                    
                    Please provide a detailed report with the following sections:
                    1. 🎯 Elevator Pitch (1-2 lines)
                    2. 👥 Target Audience & Market Size
                    3. ⚔️ Competitor Analysis (Direct & Indirect)
                    4. 💰 Monetization Strategy (How to make money)
                    5. 🚦 SWOT Analysis (Strengths, Weaknesses, Opportunities, Threats)
                    6. 🚀 30-Day Action Plan (Step-by-step roadmap to launch)
                    """

                    answer = generate_response(prompt)
                    
                    st.markdown(answer)

                else:

                    st.warning("Please enter your idea first.")

    # =====================================================
    # LOGOUT 
    # =====================================================
    
    st.sidebar.markdown("---")
    
    if st.sidebar.button("Logout"):
        
        # Clear Auto-Login Session
        if "session" in st.query_params:
            session_token = st.query_params["session"]
            cursor.execute("DELETE FROM sessions WHERE token=?", (session_token,))
            conn.commit()
            
        st.session_state.logged_in = False
        st.session_state.username = None 
            
        try:
            st.query_params.clear()
        except:
            st.experimental_set_query_params()

        st.rerun()
